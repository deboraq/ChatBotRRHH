from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


NAVY = RGBColor(20, 52, 117)
BLUE = RGBColor(45, 99, 214)
GREEN = RGBColor(21, 148, 110)
ORANGE = RGBColor(232, 147, 31)
LIGHT_BG = RGBColor(245, 248, 254)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(86, 97, 119)
DARK = RGBColor(22, 33, 58)
BORDER = RGBColor(216, 226, 243)


def style_paragraph(paragraph, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    paragraph.font.name = "Calibri"
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    paragraph.font.underline = False
    paragraph.alignment = align


def clear_text_frame(tf, word_wrap=True):
    tf.clear()
    tf.word_wrap = word_wrap


def add_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG

    header = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(0.7),
    )
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()

    brand = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(7.5), Inches(0.35))
    p = brand.text_frame.paragraphs[0]
    p.text = "BacarIT  |  Chatbot RRHH"
    style_paragraph(p, size=16, bold=True, color=WHITE)


def add_title(slide, title, subtitle="", y=1.0):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(11.8), Inches(1.4))
    tf = box.text_frame
    clear_text_frame(tf)
    p = tf.paragraphs[0]
    p.text = title
    style_paragraph(p, size=34, bold=True, color=NAVY)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        style_paragraph(p2, size=18, color=GRAY)


def add_bullet_card(slide, x, y, w, h, title, bullets, title_color=BLUE):
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BORDER

    tf = card.text_frame
    clear_text_frame(tf)
    tp = tf.paragraphs[0]
    tp.text = title
    style_paragraph(tp, size=20, bold=True, color=title_color)
    for item in bullets:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        style_paragraph(p, size=16, color=DARK)


def add_chip(slide, text, x, y, w=4.5, color=GREEN):
    chip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(0.62),
    )
    chip.fill.solid()
    chip.fill.fore_color.rgb = color
    chip.line.fill.background()
    tf = chip.text_frame
    clear_text_frame(tf)
    p = tf.paragraphs[0]
    p.text = text
    style_paragraph(p, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def add_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)

    add_title(
        slide,
        "CHATBOT RRHH",
        "Propuesta ejecutiva para mejorar servicio al colaborador y eficiencia operativa",
        y=1.45,
    )

    summary = slide.shapes.add_textbox(Inches(0.9), Inches(2.75), Inches(11.5), Inches(1.6))
    tf = summary.text_frame
    clear_text_frame(tf)
    p = tf.paragraphs[0]
    p.text = "• Atención rápida y consistente para consultas frecuentes"
    style_paragraph(p, size=19, color=DARK)
    p2 = tf.add_paragraph()
    p2.text = "• Escalamiento directo a RRHH cuando la consulta requiere intervención humana"
    style_paragraph(p2, size=19, color=DARK)
    p3 = tf.add_paragraph()
    p3.text = "• Operación por empresa/sucursal con trazabilidad completa"
    style_paragraph(p3, size=19, color=DARK)

    add_chip(slide, "Objetivo: aprobar implementación productiva", x=0.9, y=5.2, w=5.2, color=GREEN)


def add_problem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(
        slide,
        "Situación actual",
        "Hoy RRHH enfrenta demanda creciente con procesos que pueden ser más ágiles.",
    )

    add_bullet_card(
        slide,
        0.8,
        2.2,
        5.9,
        4.0,
        "Dolores principales",
        [
            "Consultas repetitivas consumen tiempo operativo.",
            "Respuesta desigual según horario o disponibilidad.",
            "Derivaciones manuales sin un circuito unificado.",
            "Menor tiempo para tareas estratégicas de RRHH.",
        ],
        title_color=ORANGE,
    )
    add_bullet_card(
        slide,
        6.7,
        2.2,
        5.8,
        4.0,
        "Oportunidad de mejora",
        [
            "Automatizar primera respuesta y clasificación.",
            "Derivar solo casos que realmente necesitan intervención humana.",
            "Medir tiempos, volumen y satisfacción en una sola vista.",
            "Escalar el servicio sin crecer linealmente en costos.",
        ],
        title_color=GREEN,
    )


def add_solution_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(
        slide,
        "La propuesta",
        "Un asistente de RRHH con atención automática + atención humana integrada.",
    )

    add_bullet_card(
        slide,
        0.8,
        2.15,
        11.7,
        1.7,
        "¿Qué hace el chatbot?",
        [
            "Responde preguntas frecuentes al instante y guía al colaborador.",
            "Si hace falta, deriva la conversación al equipo de RRHH en vivo.",
        ],
        title_color=BLUE,
    )

    add_bullet_card(
        slide,
        0.8,
        4.0,
        3.75,
        2.2,
        "Para colaboradores",
        [
            "Canal claro y rápido.",
            "Menos espera.",
            "Mejor experiencia.",
        ],
    )
    add_bullet_card(
        slide,
        4.8,
        4.0,
        3.75,
        2.2,
        "Para RRHH",
        [
            "Menos carga repetitiva.",
            "Más foco en casos críticos.",
            "Bandeja única de atención.",
        ],
    )
    add_bullet_card(
        slide,
        8.75,
        4.0,
        3.75,
        2.2,
        "Para dirección",
        [
            "Visibilidad de métricas.",
            "Escalabilidad controlada.",
            "Mejor productividad global.",
        ],
    )


def add_journey_collaborator_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Experiencia del colaborador", "Recorrido simple en 4 momentos.")

    steps = [
        "1) Consulta inicial",
        "2) Respuesta inmediata",
        "3) Si aplica: pase a RRHH",
        "4) Resolución y cierre",
    ]
    details = [
        "El colaborador escribe su consulta en lenguaje natural.",
        "El asistente responde con información clara y accionable.",
        "Cuando el caso lo requiere, un agente humano toma la conversación.",
        "Queda registro de la gestión para seguimiento y mejora.",
    ]
    x = 0.85
    for title, detail in zip(steps, details):
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(2.35),
            Inches(3.1),
            Inches(3.5),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = BORDER
        tf = box.text_frame
        clear_text_frame(tf)
        p = tf.paragraphs[0]
        p.text = title
        style_paragraph(p, size=17, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        p2 = tf.add_paragraph()
        p2.text = detail
        style_paragraph(p2, size=14, color=DARK, align=PP_ALIGN.CENTER)
        x += 3.2


def add_journey_rrhh_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Experiencia del equipo RRHH", "Gestión ordenada, colaborativa y medible.")

    add_bullet_card(
        slide,
        0.8,
        2.15,
        11.7,
        4.25,
        "Qué cambia para RRHH",
        [
            "Las conversaciones llegan organizadas y priorizadas.",
            "Se pueden tomar, reasignar y cerrar casos con trazabilidad.",
            "La operación se divide por empresa/sucursal según permisos.",
            "Se reduce la carga administrativa y mejora el tiempo de respuesta.",
            "La información queda centralizada para control y auditoría.",
        ],
        title_color=BLUE,
    )


def add_capabilities_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Capacidades de negocio ya disponibles", "")

    add_bullet_card(
        slide,
        0.8,
        2.2,
        5.75,
        4.0,
        "Operación y servicio",
        [
            "Atención inicial automática 24/7.",
            "Derivación a personas cuando el caso lo necesita.",
            "Historial de conversaciones centralizado.",
            "Dashboard con indicadores clave.",
        ],
    )
    add_bullet_card(
        slide,
        6.75,
        2.2,
        5.75,
        4.0,
        "Escala y gobierno",
        [
            "Gestión multiempresa y multisucursal.",
            "Usuarios, roles y permisos administrables.",
            "Asignación automática y manual de conversaciones.",
            "Base lista para expansión a WhatsApp.",
        ],
        title_color=GREEN,
    )


def add_needs_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(
        slide,
        "Qué vamos a necesitar para producción",
        "Componentes clave para operar con estabilidad y escalar el servicio.",
    )

    add_bullet_card(
        slide,
        0.8,
        2.2,
        3.75,
        4.0,
        "Base de datos",
        [
            "Firestore",
            "Historial y métricas",
            "Configuración multiempresa",
            "Persistencia segura",
        ],
        title_color=BLUE,
    )
    add_bullet_card(
        slide,
        4.8,
        2.2,
        3.75,
        4.0,
        "Inteligencia",
        [
            "Google AI Studio (prototipo)",
            "Gemini API para producción",
            "Mejor respuesta por contexto",
            "Escalabilidad por consumo",
        ],
        title_color=GREEN,
    )
    add_bullet_card(
        slide,
        8.8,
        2.2,
        3.75,
        4.0,
        "Canal y operación",
        [
            "WhatsApp Business API",
            "Proveedor BSP habilitado",
            "Cloud Run para backend",
            "Monitoreo y seguridad",
        ],
        title_color=ORANGE,
    )


def add_costs_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(
        slide,
        "Inversión mensual estimada (USD)",
        "Referencia para 800 a 1000 consultas mensuales. Valores orientativos sujetos a uso y país.",
    )

    scenarios = [
        (
            "Escenario base",
            "USD 120 - 220 / mes",
            ["Piloto controlado", "Volumen moderado", "Uso inicial de WhatsApp"],
            BLUE,
        ),
        (
            "Escenario objetivo",
            "USD 220 - 380 / mes",
            ["800-1000 consultas/mes", "Operación continua", "RRHH + métricas completas"],
            GREEN,
        ),
        (
            "Escenario escalado",
            "USD 380 - 650 / mes",
            ["Más empresas/sucursales", "Mayor tráfico", "Campañas y plantillas WhatsApp"],
            ORANGE,
        ),
    ]

    x = 0.8
    for name, amount, notes, color in scenarios:
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(2.2),
            Inches(3.75),
            Inches(3.8),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = BORDER
        tf = card.text_frame
        clear_text_frame(tf)
        p = tf.paragraphs[0]
        p.text = name
        style_paragraph(p, size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
        p2 = tf.add_paragraph()
        p2.text = amount
        style_paragraph(p2, size=22, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        for note in notes:
            pn = tf.add_paragraph()
            pn.text = f"• {note}"
            style_paragraph(pn, size=14, color=DARK)
        x += 4.0

    foot = slide.shapes.add_textbox(Inches(0.9), Inches(6.2), Inches(11.6), Inches(0.55))
    ft = foot.text_frame
    clear_text_frame(ft)
    fp = ft.paragraphs[0]
    fp.text = (
        "Sugerencia de presupuesto inicial al comité: USD 300/mes + 20% de contingencia "
        "durante los primeros 90 días."
    )
    style_paragraph(fp, size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


def add_roadmap_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Plan de implementación (90 días)", "Ejecución por etapas con hitos de negocio.")

    phases = [
        ("Fase 1 (0-30 días)", ["Ajustes finales del chatbot", "Validación de métricas", "Prueba controlada"]),
        ("Fase 2 (31-60 días)", ["Integración WhatsApp", "Capacitación de RRHH", "Inicio operación asistida"]),
        ("Fase 3 (61-90 días)", ["Escalado por empresa/sucursal", "Optimización de costos", "KPIs para comité"]),
    ]

    x = 0.85
    for title, lines in phases:
        add_bullet_card(slide, x, 2.2, 4.1, 3.9, title, lines, title_color=BLUE)
        x += 4.2


def add_decision_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Decisión solicitada al Comité Ejecutivo", "")

    add_bullet_card(
        slide,
        0.8,
        2.15,
        5.8,
        4.2,
        "Aprobaciones requeridas",
        [
            "Presupuesto mensual inicial para operación.",
            "Habilitación de infraestructura Google y IA.",
            "Habilitación del canal WhatsApp Business.",
            "Sponsor de negocio para seguimiento trimestral.",
        ],
        title_color=ORANGE,
    )
    add_bullet_card(
        slide,
        6.75,
        2.15,
        5.75,
        4.2,
        "Resultado esperado",
        [
            "Mejor experiencia del colaborador.",
            "Menor carga operativa del equipo RRHH.",
            "Mayor trazabilidad y control de gestión.",
            "Base escalable para nuevas unidades del grupo.",
        ],
        title_color=GREEN,
    )


def build_presentation(output_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_cover(prs)
    add_problem_slide(prs)
    add_solution_slide(prs)
    add_journey_collaborator_slide(prs)
    add_journey_rrhh_slide(prs)
    add_capabilities_slide(prs)
    add_needs_slide(prs)
    add_costs_slide(prs)
    add_roadmap_slide(prs)
    add_decision_slide(prs)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


if __name__ == "__main__":
    target = Path("docs/Presentacion_Chatbot_RRHH_Bacar.pptx")
    build_presentation(target)
    print(f"Presentación generada: {target.resolve()}")
